"""
PPTX / PPT document loader using python-pptx.

Extracts text slide-by-slide, capturing slide titles and text from
all shapes on each slide.

Note: python-pptx only supports the modern .pptx (OpenXML) format.
Legacy .ppt files are not natively supported.  The extraction service
should inform callers when a .ppt file is uploaded.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches  # noqa: F401 — imported to validate pptx install

from app.core.logging import get_logger
from app.loaders.base_loader import BaseLoader
from app.schemas.document import ExtractedPage

logger = get_logger(__name__)


class PPTXLoader(BaseLoader):
    """Extract text from PowerPoint PPTX files."""

    supported_extensions = {".pptx", ".ppt"}

    def extract(self, file_path: Path) -> list[ExtractedPage]:
        """
        Extract text from every slide of a PPTX file.

        Parameters
        ----------
        file_path : Path
            Path to the PPTX file.

        Returns
        -------
        list[ExtractedPage]
            One ExtractedPage per slide with slide_number (1-indexed)
            and slide_title where available.
        """
        self.validate_file(file_path)

        ext = file_path.suffix.lower()
        if ext == ".ppt":
            raise ValueError(
                f"Legacy .ppt format is not supported. "
                f"Please convert '{file_path.name}' to .pptx."
            )

        logger.info("PPTX extraction started: %s", file_path.name)

        pages: list[ExtractedPage] = []

        try:
            prs = Presentation(str(file_path))
        except Exception as exc:
            raise RuntimeError(
                f"Failed to open PPTX '{file_path.name}': {exc}"
            ) from exc

        if not prs.slides:
            raise ValueError(f"PPTX has no slides: {file_path.name}")

        for slide_idx, slide in enumerate(prs.slides, start=1):
            # --- Slide title ---
            slide_title = None
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                slide_title = slide.shapes.title.text_frame.text.strip() or None

            # --- Collect text from all shapes ---
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            texts.append(para_text)

                # Tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_texts.append(cell_text)
                        if row_texts:
                            texts.append(" | ".join(row_texts))

            slide_text = "\n".join(texts)

            # Skip empty slides
            if not slide_text.strip():
                logger.debug(
                    "Slide %d is blank in %s — skipping",
                    slide_idx,
                    file_path.name,
                )
                continue

            pages.append(
                ExtractedPage(
                    slide_number=slide_idx,
                    slide_title=slide_title,
                    text=slide_text,
                    char_count=len(slide_text),
                )
            )

        if not pages:
            raise ValueError(
                f"No text could be extracted from PPTX: {file_path.name}"
            )

        logger.info(
            "PPTX extraction complete: %s — %d slides extracted",
            file_path.name,
            len(pages),
        )
        return pages
