"""
Generate sample test fixture files (PDF, PPTX, DOCX) for extraction tests.

Run once:  python tests/generate_fixtures.py
"""

import sys
from pathlib import Path

# Ensure the project root is on the path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

FIXTURES_DIR = Path(__file__).resolve().parent / "fixtures"
FIXTURES_DIR.mkdir(parents=True, exist_ok=True)


def create_sample_pdf() -> None:
    """Create a 3-page sample PDF with educational content."""
    import fitz  # PyMuPDF

    doc = fitz.open()  # new empty document

    pages_content = [
        (
            "Chapter 1: Introduction to Operating Systems\n\n"
            "An operating system (OS) is system software that manages computer "
            "hardware, software resources, and provides common services for "
            "computer programs.\n\n"
            "Key Functions:\n"
            "- Process Management\n"
            "- Memory Management\n"
            "- File System Management\n"
            "- I/O Management\n"
        ),
        (
            "Chapter 2: Process Management\n\n"
            "A process is a program in execution. It includes the program code, "
            "current activity, and allocated resources.\n\n"
            "Process States:\n"
            "1. New - The process is being created.\n"
            "2. Ready - The process is waiting to be assigned to a processor.\n"
            "3. Running - Instructions are being executed.\n"
            "4. Waiting - The process is waiting for some event.\n"
            "5. Terminated - The process has finished execution.\n"
        ),
        (
            "Chapter 3: Deadlocks\n\n"
            "A deadlock is a situation where a set of processes are blocked "
            "because each process is holding a resource and waiting for "
            "another resource acquired by some other process.\n\n"
            "Banker's Algorithm:\n"
            "The Banker's Algorithm is a resource allocation and deadlock "
            "avoidance algorithm that tests for safety by simulating the "
            "allocation for predetermined maximum possible amounts of all "
            "resources.\n"
        ),
    ]

    for text in pages_content:
        page = doc.new_page(width=595, height=842)  # A4
        page.insert_text((72, 72), text, fontsize=11)

    output_path = FIXTURES_DIR / "sample.pdf"
    doc.save(str(output_path))
    doc.close()
    print(f"Created: {output_path}")


def create_sample_pptx() -> None:
    """Create a 3-slide sample PPTX with educational content."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    slides_content = [
        ("Database Fundamentals", [
            "A database is an organized collection of data.",
            "DBMS: Database Management System",
            "Key concepts: Tables, Rows, Columns, Keys",
        ]),
        ("SQL Basics", [
            "SQL: Structured Query Language",
            "DDL: CREATE, ALTER, DROP",
            "DML: SELECT, INSERT, UPDATE, DELETE",
            "DCL: GRANT, REVOKE",
        ]),
        ("Normalization", [
            "1NF: Eliminate repeating groups",
            "2NF: Remove partial dependencies",
            "3NF: Remove transitive dependencies",
            "BCNF: Every determinant is a candidate key",
        ]),
    ]

    for title, bullets in slides_content:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = bullets[0]
        for bullet in bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet

    output_path = FIXTURES_DIR / "sample.pptx"
    prs.save(str(output_path))
    print(f"Created: {output_path}")


def create_sample_docx() -> None:
    """Create a sample DOCX with headings and paragraphs."""
    from docx import Document

    doc = Document()

    doc.add_heading("Data Structures", level=1)
    doc.add_paragraph(
        "Data structures are ways of organizing and storing data so that "
        "they can be accessed and modified efficiently."
    )

    doc.add_heading("Arrays", level=2)
    doc.add_paragraph(
        "An array is a collection of items stored at contiguous memory "
        "locations. The idea is to store multiple items of the same type together."
    )

    doc.add_heading("Linked Lists", level=2)
    doc.add_paragraph(
        "A linked list is a linear data structure where each element is a "
        "separate object. Each element (node) contains a data field and a "
        "reference (link) to the next node in the sequence."
    )

    doc.add_heading("Trees", level=2)
    doc.add_paragraph(
        "A tree is a hierarchical data structure defined as a collection of nodes. "
        "Nodes represent value and nodes are connected by edges."
    )
    doc.add_paragraph(
        "Binary Tree: Each node has at most two children referred to as "
        "the left child and the right child."
    )

    output_path = FIXTURES_DIR / "sample.docx"
    doc.save(str(output_path))
    print(f"Created: {output_path}")


if __name__ == "__main__":
    create_sample_pdf()
    create_sample_pptx()
    create_sample_docx()
    print("\nAll fixtures generated successfully.")
